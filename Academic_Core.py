import os
import fitz
import sqlite3
from datetime import datetime
from dotenv import load_dotenv
from langchain_openai import ChatOpenAI
from langchain_core.messages import HumanMessage, SystemMessage
from pptx import Presentation

load_dotenv()

if not os.getenv("OPENAI_API_KEY"):
    raise ValueError("OPENAI_API_KEY не найден в .env файле")



class UserDatabase:
    def __init__(self, db_name="academic_users.db"):
        self.conn = sqlite3.connect(db_name)
        self.create_tables()

    def create_tables(self):
        cursor = self.conn.cursor()

        cursor.execute("""
        CREATE TABLE IF NOT EXISTS users (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            username TEXT UNIQUE NOT NULL,
            created_at TEXT
        )
        """)

        cursor.execute("""
        CREATE TABLE IF NOT EXISTS analyses (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            user_id INTEGER,
            file_name TEXT,
            result TEXT,
            created_at TEXT,
            FOREIGN KEY (user_id) REFERENCES users(id)
        )
        """)

        self.conn.commit()



class AcademicAssistant:
    def __init__(self):
        self.llm = ChatOpenAI(
            model="gpt-4o-mini",
            temperature=0.2
        )

    def analyze_material(self, text_content):
        template = """
        Ты — ведущий научный сотрудник. Главно отвечай на языке который задал пользователь.
        Проанализируй текст по структуре:
        1. ОСНОВНЫЕ ТЕЗИСЫ.
        2. КОНСПЕКТ ИССЛЕДОВАНИЯ, ФАЙЛА, ИЛИ МАТЕРИАЛА.
        3. МЕТОДОЛОГИЧЕСКИЙ АНАЛИЗ.
        4. РЕКОМЕНДУЕМЫЙ СПИСОК ЛИТЕРАТУРЫ.
        5. ВОЗМОЖНЫЕ ПУТИ РАЗВИТИЯ ИССЛЕДОВАНИЯ.
        6. ВОПРОСЫ ДЛЯ ОБСУЖДЕНИЯ.
        7. КРИТИКА ИЛИ ЗАМЕЧАНИЯ.


        Текст: {text}
        """

        prompt = PromptTemplate.from_template(template)
        chain = prompt | self.llm

        response = chain.invoke({"text": text_content})
        return response.content


# =========================
# PDF + PPTX BRAIN
# =========================

class AcademicPDFBrain:
    def __init__(self):
        self.llm = ChatOpenAI(
            model="gpt-4o-mini",
            temperature=0.1
        )

    def extract_text_from_pdf(self, file_path):
        text = ""
        try:
            with fitz.open(file_path) as doc:
                for page in doc:
                    text += page.get_text()
            return text
        except Exception as e:
            return f"Ошибка при чтении PDF: {e}"

    def extract_text_from_pptx(self, file_path):
        text = ""
        try:
            presentation = Presentation(file_path)
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except Exception as e:
            return f"Ошибка при чтении PPTX: {e}"

    def summarize_research(self, file_path):

        if file_path.endswith(".pdf"):
            raw_text = self.extract_text_from_pdf(file_path)

        elif file_path.endswith(".pptx"):
            raw_text = self.extract_text_from_pptx(file_path)

        else:
            return "Неподдерживаемый формат файла."

        context = raw_text[:20000]

        messages = [
            SystemMessage(
                content="Ты — научный рецензент. Проанализируй материал и выдели объект, новизну и замечания."
            ),
            HumanMessage(
                content=f"Текст работы:\n{context}"
            )
        ]

        response = self.llm.invoke(messages)
        return response.content
